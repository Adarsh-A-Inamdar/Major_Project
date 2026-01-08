import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pypdf import PdfReader

# Try to import pdf2image
try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False
except Exception:
    HAS_PDF2IMAGE = False

# --- THEME CONFIGURATION ---
THEME = {
    "bg_color": RGBColor(20, 24, 35),       # Dark Navy
    "title_color": RGBColor(255, 255, 255), # White
    "body_color": RGBColor(200, 200, 200),  # Light Grey
    "accent_color": RGBColor(0, 150, 136),  # Teal
    "font_main": "Arial"
}

def apply_slide_design(slide):
    """Applies a dark, modern background design to the slide."""
    # 1. Background Color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = THEME["bg_color"]

    # 2. Decorative Strip (Accent Color) at the botom
    left = 0
    top = Inches(7.3)
    width = Inches(10)
    height = Inches(0.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME["accent_color"]
    shape.line.fill.background() # No border

def find_pdf(filename):
    """Recursively search for the file."""
    for root, dirs, files in os.walk("."):
        if filename in files:
            return os.path.join(root, filename)
    return None

def create_presentation(filename='final.pdf', output_pptx='Project_Presentation.pptx'):
    print(f"Searching for {filename}...")
    pdf_path = find_pdf(filename)
    if not pdf_path:
        print(f"Error: Could not find {filename} in current directory or subdirectories.")
        return

    print(f"Found PDF at: {pdf_path}")
    prs = Presentation()
    
    # --------------------------
    # SLIDE 1: TITLE SLIDE
    # --------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    apply_slide_design(slide)
    
    image_inserted = False
    
    # Try placing PDF image
    if HAS_PDF2IMAGE:
        try:
            images = convert_from_path(pdf_path, first_page=1, last_page=1)
            if images:
                image_path = "temp_title.jpg"
                images[0].save(image_path, 'JPEG')
                slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                image_inserted = True
                os.remove(image_path)
        except:
            pass
            
    # Fallback Title Text
    if not image_inserted:
        try:
            reader = PdfReader(pdf_path)
            if len(reader.pages) > 0:
                text = reader.pages[0].extract_text()
                lines = [l.strip() for l in text.split('\n') if l.strip()]
                
                # Title Box
                title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                p.text = lines[0] if lines else "Project Presentation"
                p.font.size = Pt(44)
                p.font.bold = True
                p.font.color.rgb = THEME["title_color"]
                p.font.name = THEME["font_main"]
                p.alignment = PP_ALIGN.CENTER

                # Subtitle Box
                if len(lines) > 1:
                    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(3))
                    tf_sub = sub_box.text_frame
                    tf_sub.word_wrap = True
                    p_sub = tf_sub.add_paragraph()
                    p_sub.text = "\n".join(lines[1:6])
                    p_sub.font.size = Pt(24)
                    p_sub.font.color.rgb = THEME["accent_color"]
                    p_sub.font.name = THEME["font_main"]
                    p_sub.alignment = PP_ALIGN.CENTER
        except Exception as e:
            print(f"Error parsing title text: {e}")

    # --------------------------
    # CONTENT SLIDES
    # --------------------------
    reader = PdfReader(pdf_path)
    # We use Blank layout + Custom Textboxes to enforce our theme
    # (Built-in layouts often revert to white bg)
    
    for i in range(1, len(reader.pages)):
        text = reader.pages[i].extract_text()
        lines = [l.strip() for l in text.split('\n') if l.strip()]
        if not lines: continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
        apply_slide_design(slide)
        
        # 1. Title Shape
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title_shape.text_frame
        p = tf.add_paragraph()
        p.text = lines[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = THEME["title_color"]
        p.font.name = THEME["font_main"]
        
        # 2. Body Text
        if len(lines) > 1:
            body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            for line in lines[1:]:
                # Filter out junk lines
                if len(line) < 3: continue 
                
                p = tf.add_paragraph()
                p.text = "• " + line
                p.font.size = Pt(20)
                p.font.color.rgb = THEME["body_color"]
                p.font.name = THEME["font_main"]
                p.space_after = Pt(10)

    prs.save(output_pptx)
    print(f"Success! Saved to {output_pptx}")

if __name__ == "__main__":
    create_presentation()
